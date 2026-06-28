#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
╔══════════════════════════════════════════════════════════════════════════════╗
║  🤖 AI COPILOT — Assistant IA Maintenance Industrielle                       ║
║                                                                              ║
║  Page indépendante se connectant directement aux données de l'app :          ║
║    • Fichiers Excel OT et Avis (même répertoire)                             ║
║    • Historique KPI  (kpis/indicateurs_kpis.xlsx)                           ║
║                                                                              ║
║  Exécution :  streamlit run ai_copilot.py                                    ║
║  Prérequis : pip install openai python-pptx                                  ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""

# ═══════════════════════════════════════════════════════════════════════════════
#  IMPORTS
# ═══════════════════════════════════════════════════════════════════════════════
import streamlit as st
import pandas as pd
import numpy as np
import io
import os
import json
import re
import html as html_lib
from datetime import datetime
from openai import OpenAI

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    from docx.shared import Pt as DocxPt, Inches as DocxInches, Cm, RGBColor as DocxRGB
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# ═══════════════════════════════════════════════════════════════════════════════
#  CONFIGURATION DE LA PAGE
# ═══════════════════════════════════════════════════════════════════════════════
st.set_page_config(
    layout="wide",
    page_title="AI Copilot — Maintenance",
    page_icon="🤖",
    initial_sidebar_state="expanded"
)

# ═══════════════════════════════════════════════════════════════════════════════
#  CONSTANTES KPI (identiques à l'application principale)
# ═══════════════════════════════════════════════════════════════════════════════
QK = [
    "TAUX_REALISATION_CORRECTIF/PT", "OT préparation <1 mois", "OT préparation >3 mois",
    "OT préparation 1mois< <3mois", "OT planification <1 mois", "OT planification >3 mois",
    "OT planification 1mois< <3mois", "OT exécution <1 mois", "OT exécution >3 mois",
    "OT exécution 1mois< <3mois", "Performance Graissage",
    "Performance Inspection", "Performance Appels Systématiques"
]
PK = [
    "Taux d'approbation des Avis", "OT LANC ESTIME",
    "Backlog préparation caractérisé", "Backlog planification caractérisé",
    "OT CONFIME", "OT_COR_EGAL", "OT Fiabilité", "Total Avis de Panne"
]
ALL_KPI = QK + PK

CIBLE = {
    "TAUX_REALISATION_CORRECTIF/PT": 85, "OT préparation <1 mois": 80,
    "OT préparation >3 mois": 5, "OT préparation 1mois< <3mois": 15,
    "OT planification <1 mois": 80, "OT planification >3 mois": 5,
    "OT planification 1mois< <3mois": 15, "OT exécution <1 mois": 80,
    "OT exécution >3 mois": 5, "OT exécution 1mois< <3mois": 15,
    "Taux d'approbation des Avis": 95, "OT LANC ESTIME": 100,
    "Backlog préparation caractérisé": 100, "Backlog planification caractérisé": 100,
    "OT CONFIME": 100, "OT_COR_EGAL": 100, "Performance Graissage": 95,
    "Performance Inspection": 95, "Performance Appels Systématiques": 95,
    "OT Fiabilité": 100, "Total Avis de Panne": 100
}

ACT_MAP = {
    "TAUX_REALISATION_CORRECTIF/PT": "Améliorer le taux de réalisation des OT correctifs.",
    "OT préparation <1 mois": "Réduire l'âge de préparation des OT (< 1 mois).",
    "OT préparation >3 mois": "Traiter en priorité les OT avec préparation > 3 mois.",
    "OT préparation 1mois< <3mois": "Réduire les OT dont la préparation est entre 1 et 3 mois.",
    "OT planification <1 mois": "Réduire l'âge de planification des OT (< 1 mois).",
    "OT planification >3 mois": "Traiter en priorité les OT avec planification > 3 mois.",
    "OT planification 1mois< <3mois": "Réduire les OT dont la planification est entre 1 et 3 mois.",
    "OT exécution <1 mois": "Réduire l'âge d'exécution des OT (< 1 mois).",
    "OT exécution >3 mois": "Traiter en priorité les OT avec exécution > 3 mois.",
    "OT exécution 1mois< <3mois": "Réduire les OT dont l'exécution est entre 1 et 3 mois.",
    "Taux d'approbation des Avis": "Créer un OT pour chaque avis sans ordre associé.",
    "OT LANC ESTIME": "Estimer les coûts de tous les OT lancés.",
    "Backlog préparation caractérisé": "Caractériser l'intégralité du backlog de préparation.",
    "Backlog planification caractérisé": "Caractériser l'intégralité du backlog de planification.",
    "OT CONFIME": "Confirmer systématiquement les OT terminés.",
    "OT_COR_EGAL": "Rapprocher les coûts réels et les coûts budgétés.",
    "Performance Graissage": "Améliorer le taux de réalisation des OT de graissage (Type 350).",
    "Performance Inspection": "Améliorer le taux de réalisation des OT d'inspection (Types 290, 300, 310).",
    "Performance Appels Systématiques": "Améliorer le taux de réalisation des appels systématiques (Type 360).",
    "OT Fiabilité": "Maintenir la fiabilité des coûts OT à 100%.",
    "Total Avis de Panne": "Maintenir le suivi exhaustif des avis de panne."
}

KPI_RESP_MAP = {
    "TAUX_REALISATION_CORRECTIF/PT": "Chef d'atelier",
    "OT préparation <1 mois": "Préparateur BM", "OT préparation 1mois< <3mois": "Préparateur BM",
    "OT préparation >3 mois": "Préparateur BM",
    "OT planification <1 mois": "Planificateur BM", "OT planification 1mois< <3mois": "Planificateur BM",
    "OT planification >3 mois": "Planificateur BM",
    "OT exécution <1 mois": "Chef d'atelier", "OT exécution 1mois< <3mois": "Chef d'atelier",
    "OT exécution >3 mois": "Chef d'atelier",
    "Taux d'approbation des Avis": "Chef d'atelier", "OT LANC ESTIME": "Fiabilité",
    "Backlog préparation caractérisé": "Préparateur BM",
    "Backlog planification caractérisé": "Planificateur BM",
    "OT CONFIME": "Agent de saisie", "OT_COR_EGAL": "Agent de saisie",
    "Performance Graissage": "Chef d'atelier", "Performance Inspection": "Chef d'atelier",
    "Performance Appels Systématiques": "Chef d'atelier",
    "OT Fiabilité": "Fiabilité", "Total Avis de Panne": "Fiabilité"
}

LOWER_BETTER = [
    "OT préparation >3 mois", "OT planification >3 mois", "OT exécution >3 mois",
    "OT préparation 1mois< <3mois", "OT planification 1mois< <3mois",
    "OT exécution 1mois< <3mois"
]

MP_KW = ["CRPR ATPD", "CRPR ATMR", "CRPR ATER", "CRPR ATRS", "CRPR ATMO",
         "ATPD", "ATMR", "ATER", "ATRS", "ATMO"]
MPLAN_KW = ["ATPL ATEI", "ATPL ATAL", "ATPL ATER", "ATPL AGAR", "ATPL ATHS",
            "ATEI", "ATAL", "ATAS", "AGAR", "ATHS"]

QUICK_SUGGESTIONS = [
    "🔍 Quels KPI sont critiques ?",
    "⚡ Analyse les anomalies détectées",
    "📊 Compare SF1 et SF2",
    "🛠️ Propose un plan d'action",
    "📉 Pourquoi les KPI de performance sont-ils faibles ?",
    "📋 Quels sont les postes les plus problématiques ?",
    "📈 Analyse les tendances des KPI",
    "🎯 Résume la situation de la maintenance",
]


# ═══════════════════════════════════════════════════════════════════════════════
#  FONCTIONS UTILITAIRES — LECTURE & PRÉPARATION DES DONNÉES
# ═══════════════════════════════════════════════════════════════════════════════

def contient_mot(t, lm):
    """Vérifie si le texte t contient l'un des mots-clés des listes lm."""
    t = str(t)
    return any(m in t for l in lm for m in l.split())


def cat_age(a):
    """Catégorise un âge en mois."""
    if pd.isna(a):
        return "Inconnu"
    if a <= 1:
        return "<1 mois"
    elif a >= 3:
        return ">3 mois"
    return "1 mois < <3 mois"


def excr(df):
    """Exclut les lignes contenant 'cresseur' dans le poste de travail."""
    if "Poste travail princ." in df.columns:
        return df[~df["Poste travail princ."].astype(str).str.contains(
            "cresseur", case=False, na=False
        )].copy()
    return df


@st.cache_data(show_spinner=False)
def read_excel_safe(bytes_data):
    """Lit un fichier Excel en détectant automatiquement le format (.xlsx / .xls)."""
    bio = io.BytesIO(bytes_data)
    header = bytes_data[:8]

    if header[:4] in (b'PK\x03\x04', b'PK\x05\x06'):
        for engine in ['openpyxl', 'calamine']:
            try:
                return pd.read_excel(bio, engine=engine)
            except Exception:
                bio.seek(0)
                continue

    if header == b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1':
        for engine in ['xlrd', 'calamine']:
            try:
                return pd.read_excel(bio, engine=engine)
            except Exception:
                bio.seek(0)
                continue

    for engine in ['openpyxl', 'xlrd', 'calamine']:
        try:
            bio.seek(0)
            return pd.read_excel(bio, engine=engine)
        except Exception:
            continue

    raise ValueError("Format de fichier non reconnu.")


@st.cache_data(show_spinner=False)
def prepare_data(ot_bytes, av_bytes, date_str):
    """Prépare et enrichit les DataFrames OT et Avis (identique à l'app principale)."""
    raw_ot = read_excel_safe(ot_bytes)
    raw_av = read_excel_safe(av_bytes)
    raw_ot = excr(raw_ot)
    raw_av = excr(raw_av)

    # Conversion des dates
    for c in ["Créé le", "Date de début planifiée", "Date de clôture", "Début réel", "Fin réelle"]:
        if c in raw_ot.columns:
            raw_ot[c] = pd.to_datetime(raw_ot[c], errors="coerce")
    for c in ["Créé le", "Début souhaité", "Date de la clôture"]:
        if c in raw_av.columns:
            raw_av[c] = pd.to_datetime(raw_av[c], errors="coerce")

    now_ts = pd.Timestamp.today()
    df = raw_ot.copy()

    # Backlog caractérisation
    df["Backlog preparation"] = np.where(
        df["Statut utilisateur"].apply(lambda x: contient_mot(x, MP_KW)),
        "CARACTERISE", "NON CARACTERISE"
    )
    df["Backlog planification"] = np.where(
        df["Statut utilisateur"].apply(lambda x: contient_mot(x, MPLAN_KW)),
        "CARACTERISE", "NON CARACTERISE"
    )
    df["Type Carac Prep"] = df["Statut utilisateur"].apply(
        lambda x: next((kw.split()[0] for kw in MP_KW if kw in str(x)), "NON CARACTERISE")
    )
    df["Type Carac Plan"] = df["Statut utilisateur"].apply(
        lambda x: next((kw.split()[0] for kw in MPLAN_KW if kw in str(x)), "NON CARACTERISE")
    )

    # Âges en mois
    for dc, am, ac in [('Créé le', "amp", "ap"),
                        ('Date de début planifiée', "amlp", "alp"),
                        ('Date de début planifiée', "amex", "aex")]:
        if dc in df.columns:
            df[am] = ((now_ts.year - df[dc].dt.year) * 12 +
                       (now_ts.month - df[dc].dt.month)).round(2)
            df[ac] = df[am].apply(cat_age)
        else:
            df[am] = np.nan
            df[ac] = "Inconnu"

    # Flags
    df["OT CONFIME"] = np.where(
        df["Statut système"].str.contains("CLOT|TCLO", na=False) &
        df["Statut système"].str.contains("CONF", na=False), "OUI", "NON"
    )
    df["Contient SOPL"] = df["Statut utilisateur"].str.contains("SOPL", na=False).map({True: 1, False: 0})
    df["OT LANC ESTIME"] = np.where(df["Total coûts budgétés"].fillna(0) == 0, "NON", "OUI")
    df["OT_COR_EGAL"] = np.where(
        (df["Total coûts budgétés"].fillna(0) - df["Total coûts réels"].fillna(0)) == 0,
        "OUI", "NON"
    )
    df["_tw_num"] = pd.to_numeric(df.get("Type de travail", pd.Series(dtype=float)), errors="coerce")

    if "Statut système" in df.columns:
        df["Statut OT"] = df["Statut système"].fillna("").astype(str).str.strip().str.split().str[0]

    # Avis sans ordre
    avf = raw_av[
        (raw_av["Ordre"].isna() | (raw_av["Ordre"].astype(str).str.strip() == "")) &
        (raw_av["Type d'avis"].isin(["ZU", "Z4", "ZR", "ZP"]))
    ].copy()

    # Postes filtrés
    apm = sorted(df[
        df["Poste travail princ."].astype(str).str.startswith(("SF1", "SF2"), na=False)
    ]["Poste travail princ."].dropna().unique().tolist())

    return df, avf, apm, now_ts


# ═══════════════════════════════════════════════════════════════════════════════
#  FONCTIONS — CALCUL DES KPI
# ═══════════════════════════════════════════════════════════════════════════════

def calc_kpis(df_i, av_i, posts):
    """Calcule l'ensemble des KPI de performance et qualité par poste de travail."""
    res = {}
    df = df_i.copy()
    av = av_i.copy()
    res['dfp'] = df

    def ckpi(n, d, sz=100):
        return np.where(d == 0, sz, (n / d) * 100)

    def cpiv(df_sub, flag, col, p):
        return pd.pivot_table(
            df_sub[flag], index="Poste travail princ.", columns=col,
            values="Ordre", aggfunc="count", fill_value=0
        ).reindex(p, fill_value=0)

    def statut_pivot(df_sub, p, label_prefix=""):
        piv = pd.pivot_table(
            df_sub, index="Poste travail princ.", columns="Statut OT",
            values="Ordre", aggfunc="count", fill_value=0
        ).reindex(p, fill_value=0)
        for c in ["CLOT", "CRÉÉ", "LANC", "TCLO"]:
            piv[c] = piv.get(c, 0)
        piv["Realises"] = piv["CLOT"] + piv["TCLO"]
        piv["Total"] = piv[["CLOT", "CRÉÉ", "LANC", "TCLO"]].sum(axis=1)
        return piv

    # ── 1. TAUX_REALISATION_CORRECTIF/PT ──
    filt_corr = (df["Nº appel pl.entret."].fillna(0) == 0) & (df["Contient SOPL"] == 1)
    an = cpiv(df, filt_corr, "Statut OT", posts)
    for c in ["CLOT", "CRÉÉ", "LANC", "TCLO"]:
        an[c] = an.get(c, 0)
    an["OT_CLOTURES"] = an["CLOT"] + an["TCLO"]
    an["TOTAL_OT"] = an[["CLOT", "CRÉÉ", "LANC", "TCLO"]].sum(axis=1)
    an["TAUX_REALISATION_CORRECTIF/PT"] = np.where(
        an["TOTAL_OT"] == 0, 100.0, ckpi(an["OT_CLOTURES"], an["TOTAL_OT"])
    )
    res['an'] = an

    # ── 2-4. OT préparation ──
    pr = cpiv(
        df,
        (df["Statut OT"] == "CRÉÉ") &
        (df["Statut utilisateur"].str.contains(r"\bCRPR\b", case=False, na=False)),
        "ap", posts
    )
    for c in ["<1 mois", ">3 mois", "1 mois < <3 mois", "Inconnu"]:
        pr[c] = pr.get(c, 0)
    pr["Total"] = pr[["<1 mois", "1 mois < <3 mois", ">3 mois", "Inconnu"]].sum(axis=1)
    pr["OT préparation <1 mois"] = ckpi(pr["<1 mois"], pr["Total"])
    pr["OT préparation >3 mois"] = ckpi(pr[">3 mois"], pr["Total"], 0)
    pr["OT préparation 1mois< <3mois"] = ckpi(pr["1 mois < <3 mois"], pr["Total"], 0)
    res['pr'] = pr

    # ── 5-7. OT planification ──
    pl = cpiv(
        df,
        (df["Statut OT"] == "LANC") &
        (df["Statut utilisateur"].str.contains("ATPL", case=False, na=False)),
        "alp", posts
    )
    for c in ["<1 mois", ">3 mois", "1 mois < <3 mois", "Inconnu"]:
        pl[c] = pl.get(c, 0)
    pl["Total"] = pl[["<1 mois", "1 mois < <3 mois", ">3 mois", "Inconnu"]].sum(axis=1)
    pl["OT planification <1 mois"] = ckpi(pl["<1 mois"], pl["Total"])
    pl["OT planification >3 mois"] = ckpi(pl[">3 mois"], pl["Total"], 0)
    pl["OT planification 1mois< <3mois"] = ckpi(pl["1 mois < <3 mois"], pl["Total"], 0)
    res['pl'] = pl

    # ── 8-10. OT exécution ──
    ex = cpiv(
        df,
        (df["Statut OT"] == "LANC") & (df["Contient SOPL"] == 1),
        "aex", posts
    )
    for c in ["<1 mois", ">3 mois", "1 mois < <3 mois", "Inconnu"]:
        ex[c] = ex.get(c, 0)
    ex["Total"] = ex[["<1 mois", "1 mois < <3 mois", ">3 mois", "Inconnu"]].sum(axis=1)
    ex["OT exécution <1 mois"] = ckpi(ex["<1 mois"], ex["Total"])
    ex["OT exécution >3 mois"] = ckpi(ex[">3 mois"], ex["Total"], 0)
    ex["OT exécution 1mois< <3mois"] = ckpi(ex["1 mois < <3 mois"], ex["Total"], 0)
    res['ex'] = ex

    # ── 11. OT LANC ESTIME ──
    la = pd.pivot_table(
        df[df["Statut OT"] == "LANC"],
        index="Poste travail princ.", columns="OT LANC ESTIME",
        values="Ordre", aggfunc="count", fill_value=0
    ).reindex(posts, fill_value=0)
    for c in ["OUI", "NON"]:
        la[c] = la.get(c, 0)
    la["Total"] = la["OUI"] + la["NON"]
    la["OT LANC ESTIME"] = ckpi(la["OUI"], la["Total"])
    res['la'] = la

    # ── 12. Backlog préparation caractérisé ──
    pc = pd.pivot_table(
        df[df["Statut OT"] == "CRÉÉ"],
        index="Poste travail princ.", columns="Backlog preparation",
        values="Ordre", aggfunc="count", fill_value=0
    ).reindex(posts, fill_value=0)
    for c in ["CARACTERISE", "NON CARACTERISE"]:
        pc[c] = pc.get(c, 0)
    pc["Total"] = pc["CARACTERISE"] + pc["NON CARACTERISE"]
    pc["Backlog préparation caractérisé"] = ckpi(pc["CARACTERISE"], pc["Total"])
    res['pc'] = pc

    # ── 13. Backlog planification caractérisé ──
    pcl = pd.pivot_table(
        df[df["Statut OT"] == "LANC"],
        index="Poste travail princ.", columns="Backlog planification",
        values="Ordre", aggfunc="count", fill_value=0
    ).reindex(posts, fill_value=0)
    for c in ["CARACTERISE", "NON CARACTERISE"]:
        pcl[c] = pcl.get(c, 0)
    pcl["Total"] = pcl["CARACTERISE"] + pcl["NON CARACTERISE"]
    pcl["Backlog planification caractérisé"] = ckpi(pcl["CARACTERISE"], pcl["Total"])
    res['pcl'] = pcl

    # ── 14. OT CONFIME ──
    cf = pd.pivot_table(
        df, index="Poste travail princ.", columns="OT CONFIME",
        values="Ordre", aggfunc="count", fill_value=0
    ).reindex(posts, fill_value=0)
    for c in ["OUI", "NON"]:
        cf[c] = cf.get(c, 0)
    cf["Total"] = cf["OUI"] + cf["NON"]
    cf["OT CONFIME"] = ckpi(cf["OUI"], cf["Total"])
    res['cf'] = cf

    # ── 15. OT_COR_EGAL ──
    ce = pd.pivot_table(
        df, index="Poste travail princ.", columns="OT_COR_EGAL",
        values="Ordre", aggfunc="count", fill_value=0
    ).reindex(posts, fill_value=0)
    for c in ["OUI", "NON"]:
        ce[c] = ce.get(c, 0)
    ce["Total"] = ce["OUI"] + ce["NON"]
    ce["OT_COR_EGAL"] = ckpi(ce["OUI"], ce["Total"])
    res['ce'] = ce

    # ── 16. Performance Graissage (Type 350) ──
    gra = statut_pivot(df[df["_tw_num"] == 350], posts)
    gra["Performance Graissage"] = ckpi(gra["Realises"], gra["Total"])
    res['gra'] = gra

    # ── 17. Performance Inspection (Types 290, 300, 310) ──
    ins = statut_pivot(df[df["_tw_num"].isin([290, 300, 310])], posts)
    ins["Performance Inspection"] = ckpi(ins["Realises"], ins["Total"])
    res['ins'] = ins

    # ── 18. Performance Appels Systématiques (Type 360) ──
    app_sys = statut_pivot(df[df["_tw_num"] == 360], posts)
    app_sys["Performance Appels Systématiques"] = ckpi(app_sys["Realises"], app_sys["Total"])
    res['app_sys'] = app_sys

    # ── 19. OT Fiabilité ──
    clot_df = df[df["Statut OT"].isin(["CLOT", "TCLO"])]
    fiab = pd.pivot_table(
        clot_df, index="Poste travail princ.", columns="OT_COR_EGAL",
        values="Ordre", aggfunc="count", fill_value=0
    ).reindex(posts, fill_value=0)
    for c in ["OUI", "NON"]:
        fiab[c] = fiab.get(c, 0)
    fiab["Total"] = fiab["OUI"] + fiab["NON"]
    fiab["OT Fiabilité"] = ckpi(fiab["OUI"], fiab["Total"])
    res['fiab'] = fiab

    # ── 20. Taux d'approbation des Avis ──
    av_by_poste = av.groupby("Poste travail princ." if "Poste travail princ." in av.columns else av.columns[0]).size()
    av_by_poste = av_by_poste.reindex(posts, fill_value=0)
    total_av = df.groupby("Poste travail princ.").size().reindex(posts, fill_value=0)
    taux_approb = pd.Series(
        np.where((total_av + av_by_poste) == 0, 100.0, ckpi(total_av, total_av + av_by_poste)),
        index=posts, name="Taux d'approbation des Avis"
    )
    res['taux_approb'] = taux_approb
    res['av_count'] = av_by_poste

    # ── Assemblage des DataFrames de synthèse ──
    perf_data = {
        "TAUX_REALISATION_CORRECTIF/PT": res['an']["TAUX_REALISATION_CORRECTIF/PT"],
        "OT préparation <1 mois": res['pr']["OT préparation <1 mois"],
        "OT préparation >3 mois": res['pr']["OT préparation >3 mois"],
        "OT préparation 1mois< <3mois": res['pr']["OT préparation 1mois< <3mois"],
        "OT planification <1 mois": res['pl']["OT planification <1 mois"],
        "OT planification >3 mois": res['pl']["OT planification >3 mois"],
        "OT planification 1mois< <3mois": res['pl']["OT planification 1mois< <3mois"],
        "OT exécution <1 mois": res['ex']["OT exécution <1 mois"],
        "OT exécution >3 mois": res['ex']["OT exécution >3 mois"],
        "OT exécution 1mois< <3mois": res['ex']["OT exécution 1mois< <3mois"],
        "Performance Graissage": res['gra']["Performance Graissage"],
        "Performance Inspection": res['ins']["Performance Inspection"],
        "Performance Appels Systématiques": res['app_sys']["Performance Appels Systématiques"],
    }
    res['perf_df'] = pd.DataFrame(perf_data, index=posts)

    qual_data = {
        "Taux d'approbation des Avis": taux_approb,
        "OT LANC ESTIME": res['la']["OT LANC ESTIME"],
        "Backlog préparation caractérisé": res['pc']["Backlog préparation caractérisé"],
        "Backlog planification caractérisé": res['pcl']["Backlog planification caractérisé"],
        "OT CONFIME": res['cf']["OT CONFIME"],
        "OT_COR_EGAL": res['ce']["OT_COR_EGAL"],
        "OT Fiabilité": res['fiab']["OT Fiabilité"],
        "Total Avis de Panne": 100.0,
    }
    res['qual_df'] = pd.DataFrame(qual_data, index=posts)

    # ── Calcul des scores ──
    def calc_score(df_kpi, kpi_list):
        scores = []
        for poste in df_kpi.index:
            s, count = 0, 0
            for kpi in kpi_list:
                if kpi in df_kpi.columns and kpi in CIBLE:
                    val = df_kpi.loc[poste, kpi]
                    cible = CIBLE[kpi]
                    if cible > 0:
                        ratio = val / cible
                        if kpi in LOWER_BETTER:
                            s += min(ratio, 1.0) * 100 if ratio <= 1.0 else max(0, 2.0 - ratio) * 100
                        else:
                            s += min(ratio, 1.0) * 100
                    else:
                        s += 100
                    count += 1
            scores.append(s / count if count > 0 else 0)
        return pd.Series(scores, index=df_kpi.index)

    res['perf_df']['Score Performance'] = calc_score(res['perf_df'], QK)
    res['qual_df']['Score Qualite'] = calc_score(res['qual_df'], PK)

    return res


# ═══════════════════════════════════════════════════════════════════════════════
#  FONCTIONS — HISTORIQUE KPI & VARIATIONS
# ═══════════════════════════════════════════════════════════════════════════════

def load_historical_kpis(filepath):
    """Charge l'historique des KPI depuis le fichier Excel de l'application."""
    if not os.path.exists(filepath):
        return pd.DataFrame()
    try:
        from openpyxl import load_workbook
        wb = load_workbook(filepath, read_only=True, data_only=True)
    except Exception:
        return pd.DataFrame()

    records, section, headers = [], None, None
    for sheet_name in wb.sheetnames:
        try:
            ws = wb[sheet_name]
            for row in ws.iter_rows(values_only=True):
                cell0 = str(row[0]).strip() if row[0] else ""
                if "INDICATEURS DE PERFORMANCE" in cell0.upper():
                    section = "perf"; headers = None; continue
                elif "INDICATEURS DE QUALITE" in cell0.upper():
                    section = "qual"; headers = None; continue
                elif "ANOMALIES" in cell0.upper():
                    section = None; continue
                if section and headers is None and cell0:
                    headers = [str(c).strip() if c else "" for c in row]; continue
                if section and headers and cell0 and cell0 not in ("Cible", "Total general", ""):
                    entry = {"Date": sheet_name}
                    for j, h in enumerate(headers):
                        if j < len(row):
                            entry[h] = row[j]
                    entry["_section"] = section
                    records.append(entry)
        except Exception:
            continue
    wb.close()
    if not records:
        return pd.DataFrame()
    df = pd.DataFrame(records)
    df["Date_parsed"] = pd.to_datetime(df["Date"].str.replace("-", "/"), format="%d/%m/%Y", errors="coerce")
    return df.sort_values("Date_parsed").reset_index(drop=True)


def calculate_variations(hist_df):
    """Calcule les variations entre périodes successives de l'historique."""
    if hist_df.empty or "Date" not in hist_df.columns:
        return pd.DataFrame()
    dates = sorted(hist_df["Date"].unique())
    if len(dates) < 2:
        return pd.DataFrame()

    variations = []
    for i in range(1, len(dates)):
        prev_date, curr_date = dates[i - 1], dates[i]
        for sec_name, sec_val in [("perf", "Performance"), ("qual", "Qualite")]:
            prev_d = hist_df[(hist_df["Date"] == prev_date) & (hist_df["_section"] == sec_name)]
            curr_d = hist_df[(hist_df["Date"] == curr_date) & (hist_df["_section"] == sec_name)]
            if "Poste de travail" not in prev_d.columns or "Poste de travail" not in curr_d.columns:
                continue
            prev_idx = prev_d.set_index("Poste de travail")
            curr_idx = curr_d.set_index("Poste de travail")
            score_col = f"Score {sec_val}"
            if score_col not in prev_idx.columns or score_col not in curr_idx.columns:
                continue
            for poste in set(prev_idx.index) & set(curr_idx.index):
                try:
                    pv = float(prev_idx.loc[poste, score_col])
                    cv = float(curr_idx.loc[poste, score_col])
                except (ValueError, TypeError, KeyError):
                    continue
                diff = cv - pv
                pct = (diff / pv * 100) if pv != 0 else (100 if cv != 0 else 0)
                if abs(diff) <= 0.5:
                    sens = "Stable"
                elif (diff > 0 and sec_val == "Performance") or (diff > 0 and sec_val == "Qualite"):
                    sens = "Amelioration"
                else:
                    sens = "Degradation"
                variations.append({
                    "Date precedente": prev_date, "Date actuelle": curr_date,
                    "Poste": poste, "Type": sec_val,
                    "Valeur precedente": round(pv, 2), "Valeur actuelle": round(cv, 2),
                    "Ecart": round(diff, 2), "Ecart %": round(pct, 2), "Sens": sens
                })
    return pd.DataFrame(variations)


# ═══════════════════════════════════════════════════════════════════════════════
#  FONCTIONS — DÉTECTION DES FICHIERS
# ═══════════════════════════════════════════════════════════════════════════════

def detect_files():
    """Détecte automatiquement les fichiers OT/Avis dans le répertoire courant."""
    ot_patterns = ["ot.xlsx", "OT.xlsx", "ordres.xlsx", "ordres_de_travail.xlsx",
                   "ot_.xlsx", "OT_.xlsx", "donnees_ot.xlsx"]
    av_patterns = ["avis.xlsx", "AVIS.xlsx", "av.xlsx", "AV.xlsx",
                   "avis_panne.xlsx", "donnees_avis.xlsx"]

    ot_path = None
    for p in ot_patterns:
        if os.path.exists(p):
            ot_path = p
            break
    # Sinon, chercher tout .xlsx et tester les colonnes
    if ot_path is None:
        for f in os.listdir("."):
            if f.lower().endswith((".xlsx", ".xls")) and f != "indicateurs_kpis.xlsx":
                try:
                    test = pd.read_excel(f, nrows=5)
                    if "Ordre" in test.columns and "Poste travail princ." in test.columns:
                        ot_path = f
                        break
                except Exception:
                    continue

    av_path = None
    for p in av_patterns:
        if os.path.exists(p):
            av_path = p
            break
    if av_path is None:
        for f in os.listdir("."):
            if f.lower().endswith((".xlsx", ".xls")) and f != ot_path and f != "indicateurs_kpis.xlsx":
                try:
                    test = pd.read_excel(f, nrows=5)
                    if "Ordre" in test.columns and "Type d'avis" in test.columns:
                        av_path = f
                        break
                except Exception:
                    continue

    return ot_path, av_path


def get_date_from_file():
    """Lit la date depuis le fichier date.txt de l'application."""
    for path in ["date.txt", "./date.txt", "../date.txt"]:
        if os.path.exists(path):
            try:
                with open(path, "r", encoding="utf-8") as f:
                    return f.read().strip()
            except Exception:
                pass
    return pd.Timestamp.today().strftime("%d/%m/%Y")


# ═══════════════════════════════════════════════════════════════════════════════
#  FONCTIONS — CONSTRUCTION DU CONTEXTE IA
# ═══════════════════════════════════════════════════════════════════════════════

def build_ai_context(kpis, posts, date_str, hist_df, var_df):
    """Construit le contexte textuel complet des données pour l'IA."""
    ctx = []
    ctx.append(f"=== DATE DES DONNÉES : {date_str} ===\n")

    # Résumé global
    df = kpis.get('dfp', pd.DataFrame())
    ctx.append(f"Nombre total d'OT chargés : {len(df)}")
    av_count = kpis.get('av_count', pd.Series(dtype=int))
    ctx.append(f"Nombre d'avis sans ordre : {int(av_count.sum())}")
    ctx.append(f"Nombre de postes de travail analysés : {len(posts)}")
    ctx.append(f"Postes : {', '.join(posts)}\n")

    # KPI Performance
    perf_df = kpis.get('perf_df', pd.DataFrame())
    if not perf_df.empty:
        ctx.append("=== KPI DE PERFORMANCE PAR POSTE ===")
        ctx.append("(Cible indiquée entre parenthèses)")
        header = "Poste de travail | " + " | ".join(
            [f"{k} (cible:{CIBLE.get(k,'?')})" for k in QK]
        ) + " | Score Performance"
        ctx.append(header)
        ctx.append("-" * len(header))
        for poste in perf_df.index:
            vals = [f"{perf_df.loc[poste, k]:.1f}" if k in perf_df.columns else "N/A" for k in QK]
            score = f"{perf_df.loc[poste, 'Score Performance']:.1f}" if "Score Performance" in perf_df.columns else "N/A"
            ctx.append(f"{poste} | " + " | ".join(vals) + f" | {score}")
        # Moyennes
        ctx.append("-" * len(header))
        means = [f"{perf_df[k].mean():.1f}" if k in perf_df.columns else "N/A" for k in QK]
        mean_score = f"{perf_df['Score Performance'].mean():.1f}" if "Score Performance" in perf_df.columns else "N/A"
        ctx.append(f"MOYENNE | " + " | ".join(means) + f" | {mean_score}\n")

    # KPI Qualité
    qual_df = kpis.get('qual_df', pd.DataFrame())
    if not qual_df.empty:
        ctx.append("=== KPI DE QUALITÉ PAR POSTE ===")
        header = "Poste de travail | " + " | ".join(
            [f"{k} (cible:{CIBLE.get(k,'?')})" for k in PK]
        ) + " | Score Qualité"
        ctx.append(header)
        ctx.append("-" * len(header))
        for poste in qual_df.index:
            vals = [f"{qual_df.loc[poste, k]:.1f}" if k in qual_df.columns else "N/A" for k in PK]
            score = f"{qual_df.loc[poste, 'Score Qualite']:.1f}" if "Score Qualite" in qual_df.columns else "N/A"
            ctx.append(f"{poste} | " + " | ".join(vals) + f" | {score}")
        ctx.append("-" * len(header))
        means = [f"{qual_df[k].mean():.1f}" if k in qual_df.columns else "N/A" for k in PK]
        mean_score = f"{qual_df['Score Qualite'].mean():.1f}" if "Score Qualite" in qual_df.columns else "N/A"
        ctx.append(f"MOYENNE | " + " | ".join(means) + f" | {mean_score}\n")

    # Anomalies détectées
    ctx.append("=== ANOMALIES DÉTECTÉES (KPI sous cible) ===")
    anomaly_count = 0
    all_kpi_df = {**{k: perf_df for k in QK}, **{k: qual_df for k in PK}}
    for kpi_name, kpi_df in all_kpi_df.items():
        if kpi_df.empty or kpi_name not in kpi_df.columns:
            continue
        cible = CIBLE.get(kpi_name)
        if cible is None:
            continue
        for poste in kpi_df.index:
            val = kpi_df.loc[poste, kpi_name]
            is_anomaly = False
            if kpi_name in LOWER_BETTER:
                is_anomaly = val > cible
            else:
                is_anomaly = val < cible
            if is_anomaly:
                anomaly_count += 1
                resp = KPI_RESP_MAP.get(kpi_name, "N/A")
                action = ACT_MAP.get(kpi_name, "N/A")
                ecart = val - cible
                ctx.append(
                    f"• {poste} — {kpi_name} : {val:.1f}% "
                    f"(cible: {cible}%, écart: {ecart:+.1f} points) "
                    f"→ Responsable: {resp} → Action: {action}"
                )
    if anomaly_count == 0:
        ctx.append("Aucune anomalie détectée. Tous les KPI sont conformes aux cibles.")
    ctx.append(f"\nTotal anomalies : {anomaly_count}\n")

    # Statistiques OT par statut
    if not df.empty and "Statut OT" in df.columns:
        ctx.append("=== DISTRIBUTION DES OT PAR STATUT ===")
        statut_counts = df["Statut OT"].value_counts()
        for s, c in statut_counts.items():
            ctx.append(f"  {s} : {c} OT ({c / len(df) * 100:.1f}%)")
        ctx.append("")

    # Backlog
    if not df.empty and "Backlog preparation" in df.columns:
        created = df[df["Statut OT"] == "CRÉÉ"]
        if not created.empty:
            bp = created["Backlog preparation"].value_counts()
            ctx.append("=== BACKLOG PRÉPARATION ===")
            for k, v in bp.items():
                ctx.append(f"  {k} : {v} OT")
            ctx.append("")
    if not df.empty and "Backlog planification" in df.columns:
        launched = df[df["Statut OT"] == "LANC"]
        if not launched.empty:
            bl = launched["Backlog planification"].value_counts()
            ctx.append("=== BACKLOG PLANIFICATION ===")
            for k, v in bl.items():
                ctx.append(f"  {k} : {v} OT")
            ctx.append("")

    # Historique & Tendances
    if not var_df.empty:
        ctx.append("=== TENDANCES (variations entre périodes) ===")
        for _, row in var_df.iterrows():
            ctx.append(
                f"  {row['Poste']} — {row['Type']} : "
                f"{row['Valeur precedente']:.1f} → {row['Valeur actuelle']:.1f} "
                f"({row['Ecart']:+.1f} pts, {row['Ecart %']:+.1f}%) — {row['Sens']}"
            )
        ctx.append("")

    # Plans d'action par anomalie
    ctx.append("=== PLANS D'ACTION ASSOCIÉS AUX ANOMALIES ===")
    for kpi_name in ALL_KPI:
        if kpi_name in ACT_MAP:
            ctx.append(f"  • {kpi_name} → {ACT_MAP[kpi_name]}")
    ctx.append("")

    return "\n".join(ctx)


# ═══════════════════════════════════════════════════════════════════════════════
#  FONCTIONS — APPEL IA (ask_ai)
# ═══════════════════════════════════════════════════════════════════════════════

SYSTEM_PROMPT = """Tu es **AI Copilot**, un assistant IA expert en maintenance industrielle, spécialisé dans l'analyse des KPI de maintenance SAP PM.

## RÈGLES STRICTES
1. **Utilise UNIQUEMENT les données fournies** dans le contexte ci-dessous. Ne JAMAIS inventer de chiffres, de postes ou de valeurs.
2. Si une information n'est pas dans le contexte, dis-le clairement : *"Je n'ai pas cette information dans les données disponibles."*
3. Toutes tes réponses doivent être en **français**.
4. Cite **toujours les chiffres exacts** du contexte.
5. Structure tes réponses avec du **Markdown** (titres, tableaux, listes, gras).
6. Sois **précis, concis et professionnel**.
7. Quand tu identifies un problème, **propose toujours une action concrète**.

## CAPACITÉS
- Analyse des KPI de performance et de qualité
- Comparaison entre postes de travail
- Identification des anomalies (KPI sous cible)
- Analyse des causes probables
- Recommandations d'actions correctives
- Calcul de ratios et indicateurs complémentaires (MTBF, MTTR, disponibilité si les données le permettent)
- Synthèse et plan d'action structuré
- Génération de contenu pour rapports et présentations

## FORMAT DES RÉPONSES
- Utilise des tableaux Markdown pour les données chiffrées
- Utilise des émojis pour structurer visuellement (🔴🟡🟢 pour les statuts, ⚠️ pour les alertes, ✅ pour les conformités)
- Mets en **gras** les chiffres clés
- Termine par une **conclusion actionnable**

## KPI — RÈGLES D'INTERPRÉTATION
- Pour les KPI de type "taux" (ex: taux de réalisation) : **plus c'est élevé, mieux c'est**, sauf mention contraire.
- Pour les KPI "OT >3 mois" et "OT 1-3 mois" : **plus c'est bas, mieux c'est** (LOWER_BETTER).
- Un KPI est **critique** quand il est en dessous de sa cible de plus de 10 points.
- Un KPI est **en alerte** quand il est en dessous de sa cible de 5 à 10 points.
- Un KPI est **conforme** quand il atteint ou dépasse sa cible.
"""

REPORT_PROMPT = """Génère un rapport professionnel d'analyse des KPI de maintenance industrielle en te basant UNIQUEMENT sur les données du contexte.

Le rapport doit contenir EXACTEMENT les sections suivantes, chacune avec un titre Markdown de niveau ## :

## 1. Résumé Exécutif
Synthèse de 5-8 lignes de la situation globale.

## 2. Analyse des KPI de Performance
Tableau récapitulatif + analyse détaillée poste par poste. Identifier les forces et faiblesses.

## 3. Analyse des KPI de Qualité
Tableau récapitulatif + analyse détaillée.

## 4. Analyse des Anomalies
Liste exhaustive de toutes les anomalies avec sévérité (🔴 Critique, 🟡 Alerte), écart à la cible, responsable, impact.

## 5. Analyse des Causes
Pour chaque anomalie majeure, propose les causes probables (organisation, ressources, processus, outils).

## 6. Analyse des Risques
Évalue les risques associés aux anomalies identifiées (sécurité, production, coûts, réglementaire).

## 7. Recommandations
Liste numérotée de recommandations prioritaires, avec responsable et délai suggéré.

## 8. Plan d'Action
Tableau structuré : Action | Responsable | Priorité | Délai | Indicateur de suivi

## 9. Conclusion
Perspective d'amélioration et objectifs pour la prochaine période.

IMPORTANT : Utilise les chiffres exacts du contexte. Chaque affirmation doit être étayée par un chiffre."""

PPTX_PROMPT = """Génère le contenu pour une présentation PowerPoint de 8 diapositives sur l'analyse des KPI de maintenance.

Utilise EXACTEMENT ce format pour chaque diapositive (respecte les séparateurs ===SLIDE===) :

===SLIDE===
TITRE: [titre de la diapositive]
CONTENU:
[contenu texte de la diapositive, avec des listes à puces]
===SLIDE===

Les 8 diapositives doivent être :
1. TITRE : "Analyse des KPI de Maintenance" + sous-titre avec la date
2. CONTEXTE : Périmètre d'analyse, nombre d'OT, postes analysés
3. PERFORMANCE : Synthèse des KPI de performance avec les scores
4. QUALITÉ : Synthèse des KPI de qualité avec les scores
5. ANOMALIES : Liste des anomalies détectées avec sévérité
6. CAUSES : Analyse des causes racines
7. PLAN D'ACTION : Actions prioritaires structurées
8. CONCLUSION : Synthèse et prochains pas

IMPORTANT : Contenu factuel uniquement basé sur les données fournies. Utilise des listes courtes et percutantes adaptées à un format présentation."""


def ask_ai(question, context, history, api_key, base_url, model, temperature=0.3):
    """
    Fonction principale d'appel à l'IA.
    Recoit la question, le contexte données, l'historique, et les paramètres API.
    Retourne la réponse textuelle.
    """
    messages = [
        {"role": "system", "content": SYSTEM_PROMPT + "\n\n## DONNÉES DE CONTEXTE\n\n" + context}
    ]

    # Ajouter l'historique récent (derniers 10 échanges pour limiter le contexte)
    for msg in history[-10:]:
        if msg["role"] in ("user", "assistant"):
            messages.append({
                "role": msg["role"],
                "content": msg["content"]
            })

    messages.append({"role": "user", "content": question})

    try:
        client = OpenAI(api_key=api_key, base_url=base_url)
        response = client.chat.completions.create(
            model=model,
            messages=messages,
            temperature=temperature,
            max_tokens=4096
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"❌ **Erreur API** : {str(e)}\n\nVérifiez votre clé API et la configuration dans la barre latérale."


# ═══════════════════════════════════════════════════════════════════════════════
#  FONCTIONS — GÉNÉRATION DE RAPPORT
# ═══════════════════════════════════════════════════════════════════════════════

def generate_report_content(context, history, api_key, base_url, model):
    """Demande à l'IA de générer le contenu du rapport et le retourne."""
    return ask_ai(REPORT_PROMPT, context, history, api_key, base_url, model, temperature=0.2)


def create_report_docx(content, date_str):
    """Crée un document Word (.docx) à partir du contenu Markdown du rapport."""
    if not DOCX_AVAILABLE:
        return None
    try:
        doc = DocxDocument()

        # Style par défaut
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = DocxPt(11)

        lines = content.split('\n')
        for line in lines:
            stripped = line.strip()
            if not stripped:
                doc.add_paragraph('')
            elif stripped.startswith('## '):
                p = doc.add_heading(stripped[3:], level=2)
                for run in p.runs:
                    run.font.color.rgb = DocxRGB(0x1E, 0x3A, 0x5F)
            elif stripped.startswith('# '):
                p = doc.add_heading(stripped[2:], level=1)
                for run in p.runs:
                    run.font.color.rgb = DocxRGB(0x1E, 0x3A, 0x5F)
            elif stripped.startswith('### '):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith('---'):
                doc.add_paragraph('─' * 60)
            elif stripped.startswith('|'):
                # Tableau Markdown
                cells = [c.strip() for c in stripped.split('|')[1:-1]]
                p = doc.add_paragraph('  |  '.join(cells))
                p.style.font.size = DocxPt(9)
            elif re.match(r'^\d+\.', stripped):
                doc.add_paragraph(stripped, style='List Number')
            elif stripped.startswith('- ') or stripped.startswith('* '):
                doc.add_paragraph(stripped[2:], style='List Bullet')
            elif stripped.startswith('**') and stripped.endswith('**'):
                p = doc.add_paragraph()
                run = p.add_run(stripped[2:-2])
                run.bold = True
            else:
                doc.add_paragraph(stripped)

        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer
    except Exception as e:
        st.error(f"Erreur lors de la génération DOCX : {e}")
        return None


# ═══════════════════════════════════════════════════════════════════════════════
#  FONCTIONS — GÉNÉRATION POWERPOINT
# ═══════════════════════════════════════════════════════════════════════════════

def generate_pptx_content(context, history, api_key, base_url, model):
    """Demande à l'IA de générer le contenu des diapositives."""
    return ask_ai(PPTX_PROMPT, context, history, api_key, base_url, model, temperature=0.2)


def create_powerpoint(content, date_str):
    """Crée un fichier PowerPoint à partir du contenu généré par l'IA."""
    if not PPTX_AVAILABLE:
        return None
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Couleurs du thème
        PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
        ACCENT = RGBColor(0x25, 0x63, 0xEB)
        DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)

        # Parser les diapositives
        slides_data = re.split(r'===SLIDE===', content)
        slides_data = [s.strip() for s in slides_data if s.strip()]

        for idx, slide_text in enumerate(slides_data):
            # Layout blank
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)

            # Fond
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = LIGHT_BG

            # Barre latérale gauche
            left_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5)
            )
            left_bar.fill.solid()
            left_bar.fill.fore_color.rgb = PRIMARY
            left_bar.line.fill.background()

            # Barre supérieure
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = ACCENT
            top_bar.line.fill.background()

            # Parser titre et contenu
            title = ""
            body_lines = []
            for line in slide_text.split('\n'):
                line = line.strip()
                if line.upper().startswith('TITRE:'):
                    title = line[6:].strip()
                elif line.upper().startswith('CONTENU:'):
                    continue
                elif line:
                    body_lines.append(line)

            # Titre de la diapositive
            if title:
                # Fond du titre
                title_bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0.15), Inches(0.08), Inches(13.183), Inches(1.1)
                )
                title_bg.fill.solid()
                title_bg.fill.fore_color.rgb = PRIMARY
                title_bg.line.fill.background()

                txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12.5), Inches(1.0))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.LEFT

            # Contenu
            if body_lines:
                txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.5))
                tf = txBox.text_frame
                tf.word_wrap = True

                for i, line in enumerate(body_lines):
                    # Nettoyer le Markdown
                    clean = line.replace('**', '').replace('**', '').replace('*', '').replace('`', '')
                    clean = re.sub(r'#+\s*', '', clean)
                    clean = clean.strip('🔴🟡🟢⚠️✅❌📊📋🎯📈📉🔍⚡🛠️')

                    if not clean:
                        continue

                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    # Détection des listes
                    if clean.startswith('- ') or clean.startswith('• '):
                        clean = clean[2:]
                        p.level = 0
                        p.text = "▸  " + clean
                    elif re.match(r'^\d+\.', clean):
                        p.text = clean
                        p.level = 0
                    else:
                        p.text = clean
                        p.level = 0

                    p.font.size = Pt(18)
                    p.font.color.rgb = DARK_TEXT
                    p.space_after = Pt(8)

            # Numéro de diapositive
            num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.4))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(idx + 1)
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
            p.alignment = PP_ALIGN.RIGHT

        # Première diapo spéciale : titre principal
        if len(prs.slides) > 0:
            first_slide = prs.slides[0]
            # Fond gradient bleu
            bg_shape = first_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = PRIMARY
            bg_shape.line.fill.background()

            # Retirer les éléments par défaut et recréer
            for shape in list(first_slide.shapes):
                sp = shape._element
                sp.getparent().remove(sp)

            # Icône
            icon_box = first_slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(2), Inches(1.5))
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.text = "🤖"
            p.font.size = Pt(72)
            p.alignment = PP_ALIGN.CENTER

            # Titre principal
            title_box = first_slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.333), Inches(1.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = "AI Copilot — Analyse des KPI de Maintenance"
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

            # Sous-titre
            sub_box = first_slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Rapport d'analyse — {date_str}"
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
            p.alignment = PP_ALIGN.CENTER

            # Ligne décorative
            line_shape = first_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.3), Inches(4.333), Inches(0.05)
            )
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = ACCENT
            line_shape.line.fill.background()

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
    except Exception as e:
        st.error(f"Erreur lors de la génération PowerPoint : {e}")
        return None


# ═══════════════════════════════════════════════════════════════════════════════
#  CSS PERSONNALISÉE
# ═══════════════════════════════════════════════════════════════════════════════

def inject_copilot_css():
    """Injecte les styles CSS pour l'interface AI Copilot."""
    st.markdown("""<style>
    [data-testid="stHeaderActionElements"] { display: none !important; }
    [data-testid="stActionButtonContainer"] { display: none !important; }

    .stApp { background: #f0f4f8; }

    div[data-testid="stSidebar"] {
        background: linear-gradient(180deg, #1e40af 0%, #1e3a8a 50%, #1e3a5f 100%) !important;
    }
    div[data-testid="stSidebar"] * { color: rgba(255,255,255,0.9) !important; }
    div[data-testid="stSidebar"] .stSelectbox label,
    div[data-testid="stSidebar"] .stTextInput label,
    div[data-testid="stSidebar"] .stFileUploader label,
    div[data-testid="stSidebar"] .stRadio label {
        color: rgba(255,255,255,0.9) !important;
        font-weight: 600;
        font-size: 13px;
        text-transform: uppercase;
        letter-spacing: 0.5px;
    }
    div[data-testid="stSidebar"] div[data-testid="stWidget"] {
        background: rgba(255,255,255,0.1);
        border-radius: 6px;
        padding: 5px 10px;
        margin-bottom: 5px;
        border: 1px solid rgba(255,255,255,0.15);
    }
    div[data-testid="stSidebar"] .stSelectbox > div > div,
    div[data-testid="stSidebar"] .stTextInput > div > div {
        background: rgba(255,255,255,0.95) !important;
        border-radius: 5px;
    }
    div[data-testid="stSidebar"] section[data-testid="stSidebarUserContent"] {
        padding-top: 10px;
    }

    .copilot-header {
        background: linear-gradient(135deg, #1e3a5f 0%, #2563eb 60%, #3b82f6 100%);
        padding: 24px 32px;
        border-radius: 14px;
        margin-bottom: 12px;
        box-shadow: 0 8px 32px rgba(30,58,95,0.2);
        display: flex;
        align-items: center;
        gap: 16px;
    }
    .copilot-header h1 {
        color: #fff;
        font-size: 32px;
        font-weight: 800;
        margin: 0;
        flex: 1;
        font-family: 'Inter', system-ui, sans-serif;
    }
    .copilot-header .badge {
        background: rgba(255,255,255,0.2);
        padding: 6px 16px;
        border-radius: 20px;
        color: #fff;
        font-size: 13px;
        font-weight: 700;
        border: 1px solid rgba(255,255,255,0.3);
        backdrop-filter: blur(10px);
    }
    .copilot-desc {
        color: #64748b;
        font-size: 15px;
        margin-bottom: 12px;
        padding-left: 4px;
    }

    .suggestions-row {
        display: flex;
        flex-wrap: wrap;
        gap: 8px;
        margin-bottom: 12px;
    }
    .suggestion-chip {
        background: #fff;
        border: 1px solid #e2e8f0;
        border-radius: 20px;
        padding: 8px 16px;
        font-size: 13px;
        font-weight: 600;
        color: #334155;
        cursor: pointer;
        transition: all 0.2s;
        box-shadow: 0 1px 3px rgba(0,0,0,0.04);
    }
    .suggestion-chip:hover {
        background: #eff6ff;
        border-color: #3b82f6;
        color: #1e40af;
        transform: translateY(-1px);
        box-shadow: 0 4px 8px rgba(59,130,246,0.15);
    }

    .action-bar {
        display: flex;
        gap: 8px;
        margin-bottom: 12px;
        flex-wrap: wrap;
    }

    .chat-container {
        background: #fff;
        border-radius: 14px;
        border: 1px solid #e2e8f0;
        box-shadow: 0 4px 16px rgba(0,0,0,0.06);
        overflow: hidden;
    }

    /* Style des messages chat Streamlit */
    [data-testid="stChatMessage"] {
        padding: 16px 20px !important;
        font-size: 14.5px;
        line-height: 1.65;
    }
    [data-testid="stChatMessage"][data-testid*="assistant"] {
        background: #f8fafc;
    }
    [data-testid="stChatMessage"][data-testid*="user"] {
        background: linear-gradient(135deg, #eff6ff, #dbeafe);
    }
    [data-testid="stChatMessage"] p {
        margin-bottom: 8px;
    }
    [data-testid="stChatMessage"] table {
        font-size: 13px;
        width: 100%;
        border-collapse: collapse;
        margin: 10px 0;
    }
    [data-testid="stChatMessage"] table th {
        background: #1e3a5f;
        color: #fff;
        padding: 8px 12px;
        font-weight: 700;
        text-align: left;
        font-size: 12px;
    }
    [data-testid="stChatMessage"] table td {
        padding: 6px 12px;
        border-bottom: 1px solid #e2e8f0;
        font-size: 13px;
    }
    [data-testid="stChatMessage"] table tr:nth-child(even) td {
        background: #f8fafc;
    }
    [data-testid="stChatMessage"] ul, [data-testid="stChatMessage"] ol {
        padding-left: 20px;
    }
    [data-testid="stChatMessage"] li {
        margin-bottom: 4px;
    }
    [data-testid="stChatMessage"] h2, [data-testid="stChatMessage"] h3 {
        color: #1e3a5f;
        margin-top: 16px;
        margin-bottom: 8px;
    }

    .status-badge {
        display: inline-flex;
        align-items: center;
        gap: 6px;
        padding: 4px 12px;
        border-radius: 12px;
        font-size: 12px;
        font-weight: 700;
    }
    .status-ok { background: #d1fae5; color: #065f46; }
    .status-ko { background: #fee2e2; color: #991b1b; }
    .status-warn { background: #fef3c7; color: #92400e; }

    @media (max-width: 768px) {
        .copilot-header { padding: 16px; flex-direction: column; text-align: center; }
        .copilot-header h1 { font-size: 24px; }
        .suggestions-row { justify-content: center; }
    }
    </style>""", unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  INTERFACE UTILISATEUR — SIDEBAR
# ═══════════════════════════════════════════════════════════════════════════════

def render_sidebar():
    """Affiche la barre latérale avec configuration des données et de l'IA."""
    st.sidebar.markdown("## ⚙️ Configuration")

    # ── Source de données ──
    with st.sidebar.expander("📂 Sources de données", expanded=True):
        source_mode = st.radio(
            "Mode de chargement",
            ["Auto-détection", "Upload manuel"],
            label_visibility="collapsed",
            key="source_mode"
        )

        ot_file = None
        av_file = None

        if source_mode == "Auto-détection":
            ot_path, av_path = detect_files()
            if ot_path:
                st.success(f"✅ OT : `{ot_path}`")
                with open(ot_path, "rb") as f:
                    ot_file = f.read()
            else:
                st.warning("❌ Fichier OT non détecté")

            if av_path:
                st.success(f"✅ Avis : `{av_path}`")
                with open(av_path, "rb") as f:
                    av_file = f.read()
            else:
                st.warning("❌ Fichier Avis non détecté")

            if not ot_path or not av_path:
                st.info("💡 Placez les fichiers OT et Avis (.xlsx) dans le même répertoire que ce script, ou passez en mode Upload.")
        else:
            ot_up = st.file_uploader("Fichier OT (.xlsx)", type=["xlsx", "xls"], key="ot_upload")
            av_up = st.file_uploader("Fichier Avis (.xlsx)", type=["xlsx", "xls"], key="av_upload")
            if ot_up:
                ot_file = ot_up.read()
                st.success(f"✅ OT : {ot_up.name}")
            if av_up:
                av_file = av_up.read()
                st.success(f"✅ Avis : {av_up.name}")

    # ── Configuration IA ──
    with st.sidebar.expander("🔑 Configuration IA", expanded=True):
        api_key = st.text_input(
            "Clé API",
            type="password",
            value=st.session_state.get("saved_api_key", ""),
            key="api_key_input",
            help="OpenAI, OpenRouter ou toute API compatible OpenAI"
        )
        st.session_state.saved_api_key = api_key

        base_url = st.text_input(
            "URL de base API",
            value=st.session_state.get("saved_base_url", "https://api.openai.com/v1"),
            key="base_url_input",
            help="Par défaut : OpenAI. Pour OpenRouter : https://openrouter.ai/api/v1"
        )
        st.session_state.saved_base_url = base_url

        model = st.selectbox(
            "Modèle",
            [
                "gpt-4o",
                "gpt-4o-mini",
                "gpt-4-turbo",
                "gpt-3.5-turbo",
                "anthropic/claude-3.5-sonnet",
                "google/gemini-pro-1.5",
                "meta-llama/llama-3.1-70b-instruct",
            ],
            index=0,
            key="model_select"
        )

        temperature = st.slider("Température", 0.0, 1.0, 0.3, 0.1, key="temp_slider")

    # ── État des données ──
    with st.sidebar.expander("📊 État des données", expanded=False):
        if st.session_state.get("data_loaded"):
            kpis = st.session_state.kpis
            posts = st.session_state.posts
            st.success("✅ Données chargées")
            st.metric("OT analysés", len(kpis.get('dfp', pd.DataFrame())))
            av_c = kpis.get('av_count', pd.Series(dtype=int))
            st.metric("Avis sans ordre", int(av_c.sum()) if not av_c.empty else 0)
            st.metric("Postes de travail", len(posts))
            st.caption(f"Date : {st.session_state.date_str}")
        else:
            st.error("❌ Aucune donnée chargée")

        hist_path = "kpis/indicateurs_kpis.xlsx"
        if os.path.exists(hist_path):
            st.success(f"✅ Historique KPI détecté")
        else:
            st.info("ℹ️ Aucun historique KPI trouvé")

    return ot_file, av_file, api_key, base_url, model, temperature


# ═══════════════════════════════════════════════════════════════════════════════
#  INTERFACE UTILISATEUR — ZONE PRINCIPALE
# ═══════════════════════════════════════════════════════════════════════════════

def render_main():
    """Affiche l'interface principale du Copilot."""
    inject_copilot_css()

    # ── En-tête ──
    st.markdown("""
    <div class="copilot-header">
        <span style="font-size:44px">🤖</span>
        <h1>AI Copilot</h1>
        <span class="badge">Maintenance Industrielle</span>
    </div>
    """, unsafe_allow_html=True)

    st.markdown(
        '<p class="copilot-desc">💬 Posez vos questions concernant vos données maintenance. '
        "L'analyse porte sur les fichiers chargés dans la barre latérale.</p>",
        unsafe_allow_html=True
    )

    # ── Vérification données ──
    if not st.session_state.get("data_loaded"):
        st.warning("⚠️ **Aucune donnée chargée.** Veuillez configurer les sources de données dans la barre latérale puis relancer le chargement.")
        st.info("💡 **Conseil** : Placez ce fichier `ai_copilot.py` dans le même répertoire que vos fichiers Excel OT et Avis pour que l'auto-détection fonctionne.")
        return

    # ── Suggestions rapides ──
    if not st.session_state.get("copilot_messages"):
        st.markdown('<div class="suggestions-row">', unsafe_allow_html=True)
        cols = st.columns(4)
        for i, sug in enumerate(QUICK_SUGGESTIONS):
            with cols[i % 4]:
                if st.button(sug, key=f"sug_{i}", use_container_width=True):
                    st.session_state.pending_question = sug.replace(sug[:2], "").strip()
                    st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)

    # ── Barre d'actions ──
    col1, col2, col3, col4 = st.columns([1, 1, 1, 1])
    with col1:
        if st.button("🗑️ Effacer la conversation", use_container_width=True, type="secondary"):
            st.session_state.copilot_messages = []
            st.rerun()
    with col2:
        if st.button("📄 Générer un rapport", use_container_width=True, type="primary"):
            st.session_state.generate_report = True
            st.rerun()
    with col3:
        pptx_label = "📽️ Générer PowerPoint" if PPTX_AVAILABLE else "📽️ PowerPoint (lib requise)"
        pptx_disabled = not PPTX_AVAILABLE
        if st.button(pptx_label, use_container_width=True, type="primary", disabled=pptx_disabled):
            st.session_state.generate_pptx = True
            st.rerun()
    with col4:
        if st.button("📊 Résumé des données", use_container_width=True, type="secondary"):
            st.session_state.pending_question = "Donne-moi un résumé complet de la situation de la maintenance avec les chiffres clés."
            st.rerun()

    st.markdown("<br>", unsafe_allow_html=True)

    # ── Affichage des messages ──
    chat_container = st.container()

    with chat_container:
        for msg in st.session_state.copilot_messages:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"], unsafe_allow_html=False)
                # Bouton de téléchargement si présent
                if "download" in msg:
                    dl = msg["download"]
                    st.download_button(
                        label=dl["label"],
                        data=dl["data"],
                        file_name=dl["filename"],
                        mime=dl["mime"],
                        key=f"dl_{id(dl)}"
                    )

    # ── Traitement des actions spéciales ──
    api_key = st.session_state.get("saved_api_key", "")
    base_url = st.session_state.get("saved_base_url", "https://api.openai.com/v1")
    model = st.session_state.get("current_model", "gpt-4o")
    temperature = st.session_state.get("current_temperature", 0.3)
    context = st.session_state.get("ai_context", "")

    # Génération de rapport
    if st.session_state.get("generate_report"):
        st.session_state.generate_report = False
        if not api_key:
            st.error("❌ Clé API requise pour générer un rapport.")
        else:
            with st.spinner("📝 Génération du rapport en cours..."):
                st.session_state.copilot_messages.append({
                    "role": "user",
                    "content": "📄 **Génère un rapport complet d'analyse des KPI de maintenance.**"
                })
                report_content = generate_report_content(
                    context, st.session_state.copilot_messages, api_key, base_url, model
                )
                st.session_state.copilot_messages.append({
                    "role": "assistant",
                    "content": report_content
                })

                # Tentative DOCX
                docx_buffer = create_report_docx(report_content, st.session_state.date_str)
                if docx_buffer:
                    last_msg = st.session_state.copilot_messages[-1]
                    last_msg["download"] = {
                        "label": "📥 Télécharger le rapport (.docx)",
                        "data": docx_buffer,
                        "filename": f"rapport_maintenance_{st.session_state.date_str.replace('/', '-')}.docx",
                        "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    }
                else:
                    # Fallback : Markdown
                    md_bytes = report_content.encode('utf-8')
                    last_msg = st.session_state.copilot_messages[-1]
                    last_msg["download"] = {
                        "label": "📥 Télécharger le rapport (.md)",
                        "data": md_bytes,
                        "filename": f"rapport_maintenance_{st.session_state.date_str.replace('/', '-')}.md",
                        "mime": "text/markdown"
                    }
            st.rerun()

    # Génération PowerPoint
    if st.session_state.get("generate_pptx"):
        st.session_state.generate_pptx = False
        if not api_key:
            st.error("❌ Clé API requise pour générer une présentation.")
        elif not PPTX_AVAILABLE:
            st.error("❌ La librairie `python-pptx` est requise. Installez-la avec : `pip install python-pptx`")
        else:
            with st.spinner("📽️ Génération de la présentation PowerPoint..."):
                st.session_state.copilot_messages.append({
                    "role": "user",
                    "content": "📽️ **Génère une présentation PowerPoint de 8 diapositives sur l'analyse des KPI.**"
                })
                pptx_content = generate_pptx_content(
                    context, st.session_state.copilot_messages, api_key, base_url, model
                )
                pptx_buffer = create_powerpoint(pptx_content, st.session_state.date_str)
                if pptx_buffer:
                    st.session_state.copilot_messages.append({
                        "role": "assistant",
                        "content": "✅ **Présentation PowerPoint générée avec succès.** Cliquez sur le bouton ci-dessous pour la télécharger.\n\nLa présentation contient 8 diapositives couvrant le contexte, l'analyse des KPI de performance et de qualité, les anomalies, les causes, le plan d'action et la conclusion.",
                        "download": {
                            "label": "📥 Télécharger la présentation (.pptx)",
                            "data": pptx_buffer,
                            "filename": f"presentation_kpi_maintenance_{st.session_state.date_str.replace('/', '-')}.pptx",
                            "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        }
                    })
                else:
                    st.session_state.copilot_messages.append({
                        "role": "assistant",
                        "content": "❌ Erreur lors de la génération du PowerPoint. Vérifiez les logs ci-dessus."
                    })
            st.rerun()

    # ── Question en attente ──
    pending = st.session_state.get("pending_question")
    if pending:
        st.session_state.pending_question = None
        if not api_key:
            st.error("❌ Veuillez configurer votre clé API dans la barre latérale.")
        else:
            _process_question(pending, context, api_key, base_url, model, temperature)
            st.rerun()

    # ── Champ de saisie du chat ──
    if prompt := st.chat_input(
        "Posez votre question sur les données maintenance...",
        key="copilot_input",
        disabled=not st.session_state.get("data_loaded")
    ):
        if not api_key:
            st.error("❌ Veuillez configurer votre clé API dans la barre latérale.")
        else:
            _process_question(prompt, context, api_key, base_url, model, temperature)
            st.rerun()


def _process_question(question, context, api_key, base_url, model, temperature):
    """Traite une question utilisateur : l'envoie à l'IA et stocke la réponse."""
    # Ajouter la question
    st.session_state.copilot_messages.append({"role": "user", "content": question})

    # Appel IA avec spinner
    with st.spinner("🤖 Analyse en cours..."):
        response = ask_ai(
            question=question,
            context=context,
            history=st.session_state.copilot_messages,
            api_key=api_key,
            base_url=base_url,
            model=model,
            temperature=temperature
        )

    # Ajouter la réponse
    st.session_state.copilot_messages.append({"role": "assistant", "content": response})


# ═══════════════════════════════════════════════════════════════════════════════
#  POINT D'ENTRÉE PRINCIPAL
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    """Fonction principale — orchestre le chargement des données et l'affichage."""

    # ── Initialisation du session_state ──
    if "copilot_messages" not in st.session_state:
        st.session_state.copilot_messages = []
    if "data_loaded" not in st.session_state:
        st.session_state.data_loaded = False
    if "generate_report" not in st.session_state:
        st.session_state.generate_report = False
    if "generate_pptx" not in st.session_state:
        st.session_state.generate_pptx = False
    if "pending_question" not in st.session_state:
        st.session_state.pending_question = None

    # ── Sidebar ──
    ot_file, av_file, api_key, base_url, model, temperature = render_sidebar()

    # Stocker les params IA
    st.session_state.current_model = model
    st.session_state.current_temperature = temperature

    # ── Chargement des données ──
    if ot_file and av_file and not st.session_state.data_loaded:
        with st.spinner("🔄 Chargement et préparation des données..."):
            try:
                date_str = get_date_from_file()
                st.session_state.date_str = date_str

                df, avf, posts, now_ts = prepare_data(ot_file, av_file, date_str)

                if not posts:
                    st.error("❌ Aucun poste de travail SF1/SF2 détecté dans les données.")
                    return

                kpis = calc_kpis(df, avf, posts)
                st.session_state.kpis = kpis
                st.session_state.posts = posts
                st.session_state.now_ts = now_ts
                st.session_state.data_loaded = True

                # Charger l'historique KPI si disponible
                hist_path = "kpis/indicateurs_kpis.xlsx"
                if os.path.exists(hist_path):
                    hist_df = load_historical_kpis(hist_path)
                    var_df = calculate_variations(hist_df)
                    st.session_state.hist_df = hist_df
                    st.session_state.var_df = var_df
                else:
                    st.session_state.hist_df = pd.DataFrame()
                    st.session_state.var_df = pd.DataFrame()

                # Construire le contexte IA
                context = build_ai_context(
                    kpis, posts, date_str,
                    st.session_state.hist_df,
                    st.session_state.var_df
                )
                st.session_state.ai_context = context

                st.success(f"✅ {len(df)} OT et {len(avf)} avis chargés — {len(posts)} postes analysés")

            except Exception as e:
                st.error(f"❌ Erreur lors du chargement des données : {e}")
                st.exception(e)
                return

    # ── Affichage principal ──
    render_main()


# ═══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    main()
