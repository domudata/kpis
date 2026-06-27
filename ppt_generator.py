# -*- coding: utf-8 -*-
"""
Générateur de PowerPoint avancé avec graphiques et tableaux
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
import pandas as pd
import io

# Palette de couleurs
PRIMARY = RGBColor(30, 58, 95)
ACCENT = RGBColor(37, 99, 235)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(30, 41, 59)
GRAY = RGBColor(100, 116, 139)
RED = RGBColor(220, 38, 38)
GREEN = RGBColor(5, 150, 105)

def add_title(slide, text, prs_width):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs_width, Inches(0.8))
    shape.fill.solid(); shape.fill.fore_color.rgb = PRIMARY; shape.line.fill.background()
    tf = shape.text_frame; tf.margin_left = Inches(0.5)
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = WHITE

def df_to_table(slide, df, left, top, width, height):
    if df.empty: return
    rows, cols = len(df) + 1, len(df.columns)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # En-têtes
    for i, col_name in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = str(col_name)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(10); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    
    # Données
    for r_idx, row in df.iterrows():
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val) if pd.notna(val) else "0"
            p = cell.text_frame.paragraphs[0]; p.font.size = Pt(9); p.font.color.rgb = DARK; p.alignment = PP_ALIGN.CENTER

def create_pptx(json_data: dict, df_perf, df_qual, df_ano_p, df_ano_q, df_var) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    W = prs.slide_width

    # Filtrer les postes réels (exclure Total/Cible)
    if 'Poste de travail' in df_perf.columns:
        df_p = df_perf[~df_perf['Poste de travail'].astype(str).str.contains('Total|CIBLE|general', case=False, na=False)].copy()
        df_q = df_qual[~df_qual['Poste de travail'].astype(str).str.contains('Total|CIBLE|general', case=False, na=False)].copy()
    else:
        df_p, df_q = df_perf, df_qual

    postes = df_p['Poste de travail'].tolist() if not df_p.empty else []

    # ==========================================
    # SLIDE 1 : Comparaison Performance / Qualité
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "Chart Comparaison Performance / Qualité par Poste", W)
    
    if not df_p.empty and 'Score Performance' in df_p.columns and 'Score Qualite' in df_q.columns:
        chart_data = CategoryChartData()
        chart_data.categories = postes
        
        perf_vals = [float(df_p.loc[df_p['Poste de travail']==p, 'Score Performance'].iloc[0] or 0) for p in postes]
        qual_vals = [float(df_q.loc[df_q['Poste de travail']==p, 'Score Qualite'].iloc[0] or 0) for p in postes]
        
        chart_data.add_series('Performance', perf_vals)
        chart_data.add_series('Qualité', qual_vals)
        
        chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8), chart_data)
        chart = chart_frame.chart
        chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.BOTTOM; chart.legend.include_in_layout = False
        plot = chart.plots[0]; plot.gap_width = 100
    else:
        slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1)).text_frame.paragraphs[0].text = "Données insuffisantes pour générer le graphique."

    # ==========================================
    # SLIDE 2 : Taux moyens
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "Taux Moyens — Performance & Qualité", W)
    
    avg_perf = df_p['Score Performance'].mean() if 'Score Performance' in df_p.columns and not df_p.empty else 0
    avg_qual = df_q['Score Qualite'].mean() if 'Score Qualite' in df_q.columns and not df_q.empty else 0
    
    # Box Perf
    shape = slide.shapes.add_shape(1, Inches(2), Inches(2.5), Inches(4), Inches(2.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "PERFORMANCE"; p.font.size = Pt(18); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = f"{avg_perf:.1f}%"; p2.font.size = Pt(60); p.font.bold = True; p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER

    # Box Qual
    shape2 = slide.shapes.add_shape(1, Inches(7.3), Inches(2.5), Inches(4), Inches(2.5))
    shape2.fill.solid(); shape2.fill.fore_color.rgb = GREEN; shape2.line.fill.background()
    tf2 = shape2.text_frame; tf2.word_wrap = True
    p3 = tf2.paragraphs[0]; p3.text = "QUALITÉ"; p3.font.size = Pt(18); p3.font.color.rgb = WHITE; p3.alignment = PP_ALIGN.CENTER
    p4 = tf2.add_paragraph(); p4.text = f"{avg_qual:.1f}%"; p4.font.size = Pt(60); p4.font.bold = True; p4.font.color.rgb = WHITE; p4.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 3 : Détails Performance + Anomalies
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "Détails des Indicateurs de Performance", W)
    if not df_p.empty:
        df_to_table(slide, df_p, Inches(0.2), Inches(1.0), Inches(7.5), Inches(3.0))
    if not df_ano_p.empty:
        slide.shapes.add_textbox(Inches(8), Inches(1.0), Inches(5), Inches(0.5)).text_frame.paragraphs[0].text = "Anomalies à traiter :"
        df_to_table(slide, df_ano_p, Inches(8), Inches(1.5), Inches(5.1), Inches(5.0))

    # ==========================================
    # SLIDE 4 : Actions Performance
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "Actions Recommandées — Performance", W)
    actions = json_data.get("slide5", {}).get("recommendations", [])
    # Filtrer heuristiquement les actions de performance
    perf_actions = [a for a in actions if any(kw in str(a.get('action','')).lower() for kw in ['préparation', 'planification', 'exécution', 'réalisation', 'backlog', 'ot lanc'])]
    if perf_actions:
        df_to_table(slide, pd.DataFrame(perf_actions), Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5))
    else:
        slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1)).text_frame.paragraphs[0].text = "Aucune action spécifique identifiée par l'IA."

    # ==========================================
    # SLIDE 5 : Détails Qualité + Anomalies
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "Détails des Indicateurs de Qualité", W)
    if not df_q.empty:
        df_to_table(slide, df_q, Inches(0.2), Inches(1.0), Inches(7.5), Inches(3.0))
    if not df_ano_q.empty:
        slide.shapes.add_textbox(Inches(8), Inches(1.0), Inches(5), Inches(0.5)).text_frame.paragraphs[0].text = "Anomalies à traiter :"
        df_to_table(slide, df_ano_q, Inches(8), Inches(1.5), Inches(5.1), Inches(5.0))

    # ==========================================
    # SLIDE 6 : Actions Qualité
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "Actions Recommandées — Qualité", W)
    # Filtrer heuristiquement les actions de qualité
    qual_actions = [a for a in actions if any(kw in str(a.get('action','')).lower() for kw in ['graissage', 'inspection', 'appel', 'avis', 'confir', 'estim', 'coût', 'fiabilité'])]
    if qual_actions:
        df_to_table(slide, pd.DataFrame(qual_actions), Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5))
    else:
        slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1)).text_frame.paragraphs[0].text = "Aucune action spécifique identifiée par l'IA."

    # ==========================================
    # SLIDE 7 : Sparklines par Poste
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "Suivi Sparklines (Tendances) par Poste de Travail", W)
    
    if df_var is not None and not df_var.empty:
        # Créer un graphique en lignes avec les variations disponibles
        chart_data = CategoryChartData()
        dates = df_var['Date actuelle'].unique().tolist()
        if len(dates) > 0:
            chart_data.categories = dates
            for poste in postes:
                poste_data = df_var[df_var['Poste'] == poste]
                # Prendre la somme des écarts % par date pour ce poste
                values = [float(poste_data[poste_data['Date actuelle']==d]['Ecart %'].sum()) for d in dates]
                chart_data.add_series(poste, values)
            
            chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8), chart_data)
            chart = chart_frame.chart
            chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.RIGHT
        else:
            slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1)).text_frame.paragraphs[0].text = "Dates insuffisantes dans l'historique."
    else:
        slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1)).text_frame.paragraphs[0].text = "Aucune donnée d'historique (variations) disponible pour tracer les sparklines."

    # ==========================================
    # SLIDE 8 : Plan d'action global
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "📋 Plan d'Action Global", W)
    if actions:
        df_to_table(slide, pd.DataFrame(actions), Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5))
    else:
        slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1)).text_frame.paragraphs[0].text = "Aucun plan d'action généré."

    # ==========================================
    # SLIDE 9 : Conclusion
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = PRIMARY
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(1))
    tf = txBox.text_frame; p = tf.paragraphs[0]
    p.text = "Conclusion Exécutive"; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

    conclusion = json_data.get("slide6", {}).get("final_conclusion", "Aucune conclusion.")
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.333), Inches(3))
    tf2 = txBox2.text_frame; tf2.word_wrap = True; p2 = tf2.paragraphs[0]
    p2.text = f"« {conclusion} »"; p2.font.size = Pt(28); p2.font.color.rgb = WHITE; p2.font.italic = True; p2.alignment = PP_ALIGN.CENTER

    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer.getvalue()
